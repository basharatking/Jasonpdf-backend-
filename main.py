"""
JasonPDF Backend  ·  FastAPI  ·  v7.0
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Named after Jason — built with love ❤️

Endpoints:
  GET  /
  GET  /health
  POST /merge-pdf
  POST /split-pdf
  POST /compress-pdf
  POST /rotate-pdf
  POST /add-watermark
  POST /pdf-to-word
  POST /pdf-to-excel
  POST /pdf-to-jpg
  POST /jpg-to-pdf
  POST /unlock-pdf
  POST /protect-pdf
  POST /pdf-to-text
  POST /pdf-to-pptx
  POST /add-page-numbers
  POST /ocr-check
"""

import io
import os
import zipfile
import tempfile
from pathlib import Path
from typing import List

from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.responses import StreamingResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware

import pypdf
import pdfplumber
from PIL import Image
import img2pdf
from docx import Document
from docx.shared import Pt, RGBColor
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.pagesizes import A4
from reportlab.lib.utils import ImageReader

# ── PyMuPDF — safe import ─────────────────────────────────────────
try:
    import fitz
    if not hasattr(fitz, "open"):
        fitz.open = fitz.Document
    _FITZ_OK = True
except Exception:
    _FITZ_OK = False

# ── python-pptx — optional ────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt
    from pptx.dml.color import RGBColor as PptxRGB
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False

# ─── App ─────────────────────────────────────────────────────────
app = FastAPI(
    title="JasonPDF API",
    version="7.0",
    description="JasonPDF — Free online PDF tools. Named after Jason ❤️",
    docs_url="/docs",
    redoc_url="/redoc",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=[
        "X-Original-Size",
        "X-Compressed-Size",
        "X-Savings-Pct",
        "Content-Disposition",
    ],
)

# ─── Config ──────────────────────────────────────────────────────
FREE_LIMIT_MB    = int(os.environ.get("FREE_LIMIT_MB", "25"))
FREE_LIMIT_BYTES = FREE_LIMIT_MB * 1024 * 1024

TMPDIR = Path(tempfile.gettempdir()) / "jasonpdf"
TMPDIR.mkdir(exist_ok=True)

# ─── Helpers ─────────────────────────────────────────────────────
def stream_file(data: bytes, media_type: str, filename: str, extra: dict = None) -> StreamingResponse:
    headers = {
        "Content-Disposition": f'attachment; filename="{filename}"',
        "Access-Control-Expose-Headers": "X-Original-Size, X-Compressed-Size, X-Savings-Pct",
    }
    if extra:
        headers.update(extra)
    return StreamingResponse(io.BytesIO(data), media_type=media_type, headers=headers)


async def read_file(upload: UploadFile, limit_bytes: int = None) -> bytes:
    data = await upload.read()
    limit = limit_bytes or FREE_LIMIT_BYTES
    if len(data) > limit:
        raise HTTPException(413, f"File too large. Free limit is {FREE_LIMIT_MB} MB.")
    if not data:
        raise HTTPException(400, "Uploaded file is empty.")
    return data


def stem(filename: str) -> str:
    return Path(filename or "file").stem


def open_fitz(data: bytes):
    """Open PDF with fitz — works across all PyMuPDF versions."""
    if not _FITZ_OK:
        raise HTTPException(500, "PDF rendering engine not available.")
    try:
        return fitz.open(stream=data, filetype="pdf")
    except AttributeError:
        return fitz.Document(stream=data, filetype="pdf")
    except Exception as e:
        raise HTTPException(400, f"Could not open PDF: {e}")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# INFO
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

@app.get("/", tags=["Info"])
def root():
    return {
        "name": "JasonPDF API",
        "version": "7.0",
        "status": "running",
        "named_after": "Jason ❤️",
        "fitz_available": _FITZ_OK,
        "pptx_available": _PPTX_OK,
        "docs": "/docs",
        "tools": 16,
    }


@app.get("/health", tags=["Info"])
def health():
    return {
        "status": "ok",
        "version": "7.0",
        "fitz": _FITZ_OK,
        "pptx": _PPTX_OK,
        "limit_mb": FREE_LIMIT_MB,
    }


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TOOL ENDPOINTS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

# ── 1. MERGE PDF ─────────────────────────────────────────────────
@app.post("/merge-pdf", tags=["Tools"])
async def merge_pdf(files: List[UploadFile] = File(...)):
    if len(files) < 2:
        raise HTTPException(400, "Please provide at least 2 PDF files to merge.")

    writer = pypdf.PdfWriter()
    for upload in files:
        data = await read_file(upload)
        try:
            reader = pypdf.PdfReader(io.BytesIO(data))
            for page in reader.pages:
                writer.add_page(page)
        except Exception as e:
            raise HTTPException(400, f"Could not read '{upload.filename}': {e}")

    out = io.BytesIO()
    writer.write(out)
    return stream_file(out.getvalue(), "application/pdf", "merged.pdf")


# ── 2. SPLIT PDF ─────────────────────────────────────────────────
@app.post("/split-pdf", tags=["Tools"])
async def split_pdf(
    file: UploadFile = File(...),
    mode: str = Form("each"),
    start_page: int = Form(1),
    end_page: int = Form(1),
):
    data = await read_file(file)
    try:
        reader = pypdf.PdfReader(io.BytesIO(data))
        total = len(reader.pages)
    except Exception as e:
        raise HTTPException(400, f"Could not read PDF: {e}")

    if mode == "range":
        s = max(1, start_page) - 1
        e = min(total, end_page)
        if s >= e:
            raise HTTPException(400, f"Invalid page range. PDF has {total} pages.")
        writer = pypdf.PdfWriter()
        for i in range(s, e):
            writer.add_page(reader.pages[i])
        out = io.BytesIO()
        writer.write(out)
        return stream_file(out.getvalue(), "application/pdf",
                           f"{stem(file.filename)}_pages_{s+1}-{e}.pdf")

    # Split each page into individual PDFs → ZIP
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for i in range(total):
            w = pypdf.PdfWriter()
            w.add_page(reader.pages[i])
            pb = io.BytesIO()
            w.write(pb)
            zf.writestr(f"page_{str(i+1).zfill(3)}.pdf", pb.getvalue())
    return stream_file(zip_buf.getvalue(), "application/zip", "split_pages.zip")


# ── 3. COMPRESS PDF ──────────────────────────────────────────────
@app.post("/compress-pdf", tags=["Tools"])
async def compress_pdf(
    file: UploadFile = File(...),
    level: str = Form("medium"),
):
    data = await read_file(file)
    orig_size = len(data)

    quality_map = {"low": 85, "medium": 60, "high": 35}
    max_dim_map = {"low": 1600, "medium": 1200, "high": 900}
    q       = quality_map.get(level, 60)
    max_dim = max_dim_map.get(level, 1200)

    try:
        doc = open_fitz(data)
        for page in doc:
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    base = doc.extract_image(xref)
                    pil  = Image.open(io.BytesIO(base["image"])).convert("RGB")
                    w, h = pil.size
                    if max(w, h) > max_dim:
                        scale = max_dim / max(w, h)
                        pil   = pil.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
                    buf = io.BytesIO()
                    pil.save(buf, format="JPEG", quality=q, optimize=True)
                    doc.update_stream(xref, buf.getvalue())
                except Exception:
                    pass

        out = io.BytesIO()
        doc.save(out, garbage=4, deflate=True, clean=True)
        doc.close()
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Compression failed: {e}")

    comp      = out.getvalue()
    comp_size = len(comp)
    savings   = round((1 - comp_size / orig_size) * 100, 1) if orig_size else 0

    return stream_file(
        comp, "application/pdf", f"compressed_{file.filename}",
        {
            "X-Original-Size":   str(orig_size),
            "X-Compressed-Size": str(comp_size),
            "X-Savings-Pct":     str(savings),
        },
    )


# ── 4. ROTATE PDF ────────────────────────────────────────────────
@app.post("/rotate-pdf", tags=["Tools"])
async def rotate_pdf(
    file: UploadFile = File(...),
    angle: int = Form(90),
    pages: str = Form("all"),
):
    data = await read_file(file)
    if angle not in (90, 180, 270):
        raise HTTPException(400, "Angle must be 90, 180, or 270.")

    try:
        doc = open_fitz(data)
        for i, page in enumerate(doc):
            pn = i + 1
            if pages == "odd"  and pn % 2 == 0: continue
            if pages == "even" and pn % 2 != 0: continue
            page.set_rotation((page.rotation + angle) % 360)
        out = io.BytesIO()
        doc.save(out)
        doc.close()
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Rotation failed: {e}")

    return stream_file(out.getvalue(), "application/pdf", f"rotated_{file.filename}")


# ── 5. WATERMARK PDF ─────────────────────────────────────────────
@app.post("/add-watermark", tags=["Tools"])
async def add_watermark(
    file: UploadFile = File(...),
    text: str = Form("CONFIDENTIAL"),
    opacity: float = Form(0.2),
    position: str = Form("center"),
):
    if not text.strip():
        raise HTTPException(400, "Watermark text cannot be empty.")
    opacity = max(0.05, min(opacity, 0.95))

    data = await read_file(file)
    try:
        doc = open_fitz(data)
        for page in doc:
            w, h  = page.rect.width, page.rect.height
            fs    = min(w, h) * 0.08
            color = (0.55, 0.55, 0.55)
            if position == "center":
                page.insert_text(
                    fitz.Point(w * 0.15, h * 0.55), text,
                    fontsize=fs, rotate=45,
                    color=color, fill_opacity=opacity, overlay=True,
                )
            elif position == "top":
                page.insert_text(
                    fitz.Point(w * 0.5 - len(text) * fs * 0.25, h - fs - 20), text,
                    fontsize=fs, color=color, fill_opacity=opacity, overlay=True,
                )
            else:  # bottom
                page.insert_text(
                    fitz.Point(w * 0.5 - len(text) * fs * 0.25, fs + 20), text,
                    fontsize=fs, color=color, fill_opacity=opacity, overlay=True,
                )
        out = io.BytesIO()
        doc.save(out)
        doc.close()
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Watermark failed: {e}")

    return stream_file(out.getvalue(), "application/pdf", f"watermarked_{file.filename}")


# ── 6. PDF TO WORD ───────────────────────────────────────────────
@app.post("/pdf-to-word", tags=["Tools"])
async def pdf_to_word(file: UploadFile = File(...)):
    data = await read_file(file)
    doc  = Document()
    doc.styles["Normal"].font.name = "Calibri"
    doc.styles["Normal"].font.size = Pt(11)

    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            total = len(pdf.pages)
            if total == 0:
                raise HTTPException(400, "PDF has no pages.")
            for i, page in enumerate(pdf.pages):
                h = doc.add_heading(f"Page {i + 1}", level=1)
                h.runs[0].font.color.rgb = RGBColor(0x25, 0x63, 0xEB)
                text = page.extract_text() or ""
                if text.strip():
                    for line in text.split("\n"):
                        p = doc.add_paragraph(line)
                        p.paragraph_format.space_after = Pt(2)
                else:
                    p = doc.add_paragraph("[No extractable text — may be a scanned PDF]")
                    p.runs[0].italic = True
                if i < total - 1:
                    doc.add_page_break()
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Conversion failed: {e}")

    out = io.BytesIO()
    doc.save(out)
    return stream_file(
        out.getvalue(),
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        f"{stem(file.filename)}.docx",
    )


# ── 7. PDF TO EXCEL (Smart Multi-Sheet) ──────────────────────────
@app.post("/pdf-to-excel", tags=["Tools"])
async def pdf_to_excel(
    file: UploadFile = File(...),
    mode: str = Form("smart"),
):
    data = await read_file(file)
    wb   = openpyxl.Workbook()
    wb.remove(wb.active)

    hdr_fill  = PatternFill("solid", fgColor="2563EB")
    hdr_font  = Font(bold=True, color="FFFFFF", size=11)
    hdr_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    ts        = Side(style="thin", color="DDDDDD")
    tbdr      = Border(left=ts, right=ts, top=ts, bottom=ts)
    alt_fill  = PatternFill("solid", fgColor="EFF6FF")

    table_groups: dict = {}
    text_lines: list   = []

    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                if mode in ("smart", "tables"):
                    tables = page.extract_tables() or []
                    for tbl in tables:
                        if not tbl or len(tbl) < 1:
                            continue
                        raw_hdrs = tuple(str(c or "").strip() for c in tbl[0])
                        if not any(raw_hdrs):
                            continue
                        if raw_hdrs not in table_groups:
                            table_groups[raw_hdrs] = []
                        for row in tbl[1:]:
                            if any(v for v in row):
                                table_groups[raw_hdrs].append(
                                    [str(v or "").strip() for v in row]
                                )

                if mode == "text" or (mode == "smart" and not table_groups):
                    text = page.extract_text() or ""
                    for ln in text.split("\n"):
                        if ln.strip():
                            text_lines.append((page_num, ln.strip()))
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Excel conversion failed: {e}")

    if table_groups:
        for sheet_idx, (hdrs, rows) in enumerate(table_groups.items(), start=1):
            sheet_name = " & ".join(h for h in hdrs[:2] if h)[:28] or f"Table {sheet_idx}"
            existing = [ws.title for ws in wb.worksheets]
            base, count = sheet_name, 1
            while sheet_name in existing:
                sheet_name = f"{base[:25]}_{count}"
                count += 1

            ws = wb.create_sheet(title=sheet_name)
            ws.freeze_panes = "A2"

            for ci, h in enumerate(hdrs, start=1):
                cell = ws.cell(row=1, column=ci, value=h)
                cell.fill = hdr_fill; cell.font = hdr_font
                cell.alignment = hdr_align; cell.border = tbdr

            for ri, row in enumerate(rows, start=2):
                fill = alt_fill if ri % 2 == 0 else None
                for ci, val in enumerate(row, start=1):
                    c = ws.cell(row=ri, column=ci, value=val)
                    c.border = tbdr
                    if fill: c.fill = fill

            for col in ws.columns:
                max_len = max((len(str(c.value or "")) for c in col), default=10)
                ws.column_dimensions[col[0].column_letter].width = min(max(max_len + 3, 12), 55)

    elif text_lines:
        ws = wb.create_sheet(title="Extracted Text")
        ws.freeze_panes = "A2"
        for ci, h in enumerate(["Page", "Text"], start=1):
            c = ws.cell(row=1, column=ci, value=h)
            c.fill = hdr_fill; c.font = hdr_font
            c.alignment = hdr_align; c.border = tbdr
        for ri, (pg, line) in enumerate(text_lines, start=2):
            ws.cell(row=ri, column=1, value=pg).border = tbdr
            ws.cell(row=ri, column=2, value=line).border = tbdr
        ws.column_dimensions["A"].width = 8
        ws.column_dimensions["B"].width = 80
    else:
        ws = wb.create_sheet(title="No Data")
        ws.cell(row=1, column=1, value="No tables or text found. PDF may be scanned/image-based.")

    out = io.BytesIO()
    wb.save(out)
    return stream_file(
        out.getvalue(),
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        f"{stem(file.filename)}.xlsx",
    )


# ── 8. PDF TO JPG ────────────────────────────────────────────────
@app.post("/pdf-to-jpg", tags=["Tools"])
async def pdf_to_jpg(
    file: UploadFile = File(...),
    dpi: int = Form(150),
):
    data = await read_file(file)
    dpi  = max(72, min(dpi, 300))

    try:
        doc     = open_fitz(data)
        mat     = fitz.Matrix(dpi / 72.0, dpi / 72.0)
        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=mat, alpha=False)
                zf.writestr(f"page_{str(i+1).zfill(3)}.jpg", pix.tobytes("jpeg"))
        doc.close()
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Conversion failed: {e}")

    return stream_file(zip_buf.getvalue(), "application/zip",
                       f"{stem(file.filename)}_images.zip")


# ── 9. IMAGE → PDF ───────────────────────────────────────────────
@app.post("/jpg-to-pdf", tags=["Tools"])
async def jpg_to_pdf(files: List[UploadFile] = File(...)):
    if not files:
        raise HTTPException(400, "No files provided.")

    images    = []
    txt_pages = []

    for upload in files:
        raw = await read_file(upload)
        ext = Path(upload.filename or "").suffix.lower()
        if ext in {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif"}:
            img = Image.open(io.BytesIO(raw)).convert("RGB")
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=92)
            images.append(buf.getvalue())
        elif ext == ".txt":
            txt_pages.append(raw.decode("utf-8", errors="replace"))
        else:
            raise HTTPException(400, f"Unsupported file type: {upload.filename}")

    out = io.BytesIO()
    try:
        if images and not txt_pages:
            out.write(img2pdf.convert(images))
        else:
            c    = rl_canvas.Canvas(out, pagesize=A4)
            W, H = A4
            for ib in images:
                pi     = Image.open(io.BytesIO(ib))
                iw, ih = pi.size
                r      = min(W / iw, H / ih, 1.0)
                dw, dh = iw * r, ih * r
                c.drawImage(ImageReader(io.BytesIO(ib)), (W-dw)/2, (H-dh)/2, dw, dh)
                c.showPage()
            for txt in txt_pages:
                c.setFont("Helvetica", 11)
                m, lh, y = 50, 16, H - 50
                for line in txt.split("\n"):
                    if y < 66:
                        c.showPage()
                        c.setFont("Helvetica", 11)
                        y = H - 50
                    c.drawString(m, y, line[:110])
                    y -= lh
                c.showPage()
            c.save()
    except Exception as e:
        raise HTTPException(500, f"PDF creation failed: {e}")

    return stream_file(out.getvalue(), "application/pdf", "converted.pdf")


# ── 10. UNLOCK PDF ───────────────────────────────────────────────
@app.post("/unlock-pdf", tags=["Tools"])
async def unlock_pdf(
    file: UploadFile = File(...),
    password: str = Form(""),
):
    data = await read_file(file)
    try:
        doc = open_fitz(data)
        if doc.is_encrypted:
            if not doc.authenticate(password):
                raise HTTPException(400, "Incorrect password. Could not unlock PDF.")
        out = io.BytesIO()
        doc.save(out, encryption=fitz.PDF_ENCRYPT_NONE)
        doc.close()
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Unlock failed: {e}")

    return stream_file(out.getvalue(), "application/pdf", f"unlocked_{file.filename}")


# ── 11. PROTECT PDF ──────────────────────────────────────────────
@app.post("/protect-pdf", tags=["Tools"])
async def protect_pdf(
    file: UploadFile = File(...),
    password: str = Form(...),
):
    if not password:
        raise HTTPException(400, "Password cannot be empty.")

    data = await read_file(file)
    try:
        doc = open_fitz(data)
        out = io.BytesIO()
        doc.save(
            out,
            encryption=fitz.PDF_ENCRYPT_AES_256,
            user_pw=password,
            owner_pw=password + "_owner",
        )
        doc.close()
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Encryption failed: {e}")

    return stream_file(out.getvalue(), "application/pdf", f"protected_{file.filename}")


# ── 12. PDF TO TEXT ──────────────────────────────────────────────
@app.post("/pdf-to-text", tags=["Tools"])
async def pdf_to_text(file: UploadFile = File(...)):
    data  = await read_file(file)
    lines = []
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            total = len(pdf.pages)
            if total == 0:
                raise HTTPException(400, "PDF has no pages.")
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if text.strip():
                    lines.append(f"=== Page {i+1} ===\n{text.strip()}")
                else:
                    lines.append(f"=== Page {i+1} ===\n[No extractable text — may be scanned]")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Text extraction failed: {e}")

    txt_bytes = "\n\n".join(lines).encode("utf-8")
    return stream_file(txt_bytes, "text/plain", f"{stem(file.filename)}.txt")


# ── 13. PDF TO PPTX ──────────────────────────────────────────────
@app.post("/pdf-to-pptx", tags=["Tools"])
async def pdf_to_pptx(file: UploadFile = File(...)):
    if not _PPTX_OK:
        raise HTTPException(500, "python-pptx not installed. Add it to requirements.txt.")

    data = await read_file(file)
    prs  = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout     = prs.slide_layouts[6]

    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            total = len(pdf.pages)
            if total == 0:
                raise HTTPException(400, "PDF has no pages.")

            for i, page in enumerate(pdf.pages):
                slide = prs.slides.add_slide(blank_layout)

                # Page number badge
                pg_box = slide.shapes.add_textbox(Inches(12.2), Inches(0.1), Inches(1.0), Inches(0.3))
                pg_tf  = pg_box.text_frame
                pg_tf.text = f"{i+1}/{total}"
                pg_tf.paragraphs[0].runs[0].font.size = PptPt(9)
                pg_tf.paragraphs[0].runs[0].font.color.rgb = PptxRGB(0xBB, 0xBB, 0xBB)

                text = page.extract_text() or ""
                if text.strip():
                    lines = [l.strip() for l in text.split("\n") if l.strip()]
                    title_text = lines[0][:100] if lines else f"Page {i+1}"
                    body_text  = "\n".join(lines[1:]) if len(lines) > 1 else ""

                    # Title
                    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(1.0))
                    tf = title_box.text_frame
                    tf.word_wrap = True
                    p  = tf.paragraphs[0]
                    p.text = title_text
                    run = p.runs[0]
                    run.font.size = PptPt(24)
                    run.font.bold = True
                    run.font.color.rgb = PptxRGB(0x1A, 0x1A, 0x1A)

                    # Body
                    if body_text:
                        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.5))
                        btf = body_box.text_frame
                        btf.word_wrap = True
                        bp  = btf.paragraphs[0]
                        bp.text = body_text[:2000]
                        bp.runs[0].font.size = PptPt(14)
                        bp.runs[0].font.color.rgb = PptxRGB(0x3A, 0x3A, 0x3A)
                else:
                    empty_box = slide.shapes.add_textbox(Inches(4), Inches(3), Inches(5), Inches(1.5))
                    etf = empty_box.text_frame
                    ep  = etf.paragraphs[0]
                    ep.text = f"Page {i+1}"
                    ep.runs[0].font.size = PptPt(32)
                    ep.runs[0].font.color.rgb = PptxRGB(0xCC, 0xCC, 0xCC)

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"PPTX conversion failed: {e}")

    out = io.BytesIO()
    prs.save(out)
    return stream_file(
        out.getvalue(),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        f"{stem(file.filename)}.pptx",
    )


# ── 14. ADD PAGE NUMBERS ─────────────────────────────────────────
@app.post("/add-page-numbers", tags=["Tools"])
async def add_page_numbers(
    file: UploadFile = File(...),
    position: str = Form("bottom-center"),
    format: str   = Form("number"),
    start: int    = Form(1),
):
    data = await read_file(file)
    try:
        doc   = open_fitz(data)
        total = len(doc)
        for i, page in enumerate(doc):
            num = start + i
            w, h = page.rect.width, page.rect.height
            fs   = 10

            if format == "page_of":
                label = f"Page {num} of {total}"
            elif format == "dashes":
                label = f"— {num} —"
            else:
                label = str(num)

            margin = 20
            tw = len(label) * fs * 0.55

            pos_map = {
                "bottom-center": fitz.Point(w / 2 - tw / 2, margin),
                "bottom-right":  fitz.Point(w - tw - margin, margin),
                "bottom-left":   fitz.Point(margin, margin),
                "top-center":    fitz.Point(w / 2 - tw / 2, h - margin - fs),
                "top-right":     fitz.Point(w - tw - margin, h - margin - fs),
                "top-left":      fitz.Point(margin, h - margin - fs),
            }
            pt = pos_map.get(position, fitz.Point(w / 2 - tw / 2, margin))

            page.insert_text(pt, label, fontsize=fs, color=(0.3, 0.3, 0.3), overlay=True)

        out = io.BytesIO()
        doc.save(out)
        doc.close()
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Page numbering failed: {e}")

    return stream_file(out.getvalue(), "application/pdf", f"numbered_{file.filename}")


# ── 15. OCR CHECK ────────────────────────────────────────────────
@app.post("/ocr-check", tags=["Tools"])
async def ocr_check(file: UploadFile = File(...)):
    data        = await read_file(file)
    total_chars = 0
    total_pages = 0

    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            total_pages = len(pdf.pages)
            for page in pdf.pages[:5]:
                total_chars += len((page.extract_text() or "").strip())
    except Exception as e:
        raise HTTPException(400, f"Could not read PDF: {e}")

    avg        = total_chars / max(total_pages, 1)
    is_scanned = avg < 50

    return JSONResponse({
        "is_scanned":         is_scanned,
        "avg_chars_per_page": round(avg, 1),
        "total_pages":        total_pages,
        "message": (
            "Scanned PDF detected. Text extraction may not work."
            if is_scanned else
            "PDF has a text layer and should convert accurately."
        ),
    })
